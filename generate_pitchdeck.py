import os
from pptx import Presentation
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()


def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = f"{os.getenv('WEBSITE_NAME')}\n"
    subtitle.text = "Your Gateway to Cryptocurrency Intelligence"


def create_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "The Problem"
    content.text = "Cryptocurrency market information is fragmented\n"
    content.text += "Lack of reliable, real-time crypto data\n"
    content.text += "Complex technical information barrier for new users\n"
    content.text += "Need for educational resources in crypto space"


def create_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Our Solution"
    content.text = "Real-time cryptocurrency price tracking\n"
    content.text += "Comprehensive news aggregation\n"
    content.text += "Educational resources and whitepapers\n"
    content.text += "User-friendly interface for all experience levels"


def create_features_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Key Features"
    content.text = "Live cryptocurrency price updates\n"
    content.text += "RSS feed integration\n"
    content.text += "Cryptocurrency aggregated analysis\n"


def create_values_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Our Values"
    content.text = "Accuracy: Providing verified, reliable information\n"
    content.text += "Transparency: Open and honest reporting\n"
    content.text += "Education: Focus on user learning\n"
    content.text += "Innovation: Leading in crypto technology"


def create_tech_stack_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Technical Stack"
    content.text = "Python/Django Backend\n"
    content.text += "Bootstrap Frontend\n"
    content.text += "SQLite Database\n"
    content.text += "RESTful API Integration\n"
    content.text += "Responsive Web Design"


def create_roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Roadmap"
    content.text = "2025:\n"
    content.text += "Lunach mobile app\n"
    content.text += "Integrate more crypto data sources\n\n"
    content.text += "Q1 2026:\n"
    content.text += "Add portfolio tracking\n"
    content.text += "Implement AI-powered insights\n"
    content.text += "Train personalize bot and agent"


def create_contact_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Contact Us"
    content.text = f"Website: {os.getenv('WEBSITE_DOMAIN')}\n"
    content.text += f"Email: hi@{os.getenv('WEBSITE_DOMAIN')}\n"
    content.text += f"Twitter: {os.getenv('WEBSITE_TWITTER')}\n"
    content.text += f"LinkedIn: {os.getenv('WEBSITE_LINKEDIN')}"


def create_investment_plan_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Investment Plan"
    content.text = "Phase 1: Seed Funding\n"
    content.text += "Objective: Develop MVP\n"
    content.text += "Amount: $50,000\n\n"
    content.text += "Phase 2: Series A\n"
    content.text += "Objective: Market Expansion\n"
    content.text += "Amount: $500,000\n"


def main():
    prs = Presentation()

    # Create all slides
    create_title_slide(prs)
    create_problem_slide(prs)
    create_solution_slide(prs)
    create_features_slide(prs)
    create_values_slide(prs)
    create_tech_stack_slide(prs)
    create_roadmap_slide(prs)
    create_investment_plan_slide(prs)
    create_contact_slide(prs)

    # Save the presentation
    prs.save("pitchdeck.pptx")
    print("Pitch deck generated successfully!")


if __name__ == "__main__":
    main()
